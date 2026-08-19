"""Build paired A2-NLP poster concepts from the UW templates.

Outputs are editable PowerPoint files plus print PDFs and PNG previews. The layout functions use
recursive/axis-aligned splits: each design starts with one content rectangle and partitions it
into a small number of visual regions before placing text or figures.

FROM PR #77 VIA #78. Leon wrote this and it is the only tool in the project that turns a board
into something printable. Four notes for whoever picks it up.

`build_board` IS THE POSTER. It sets report 12's measured grid -- nine cells in a 3x3 and three
strip blocks across the foot -- from whichever board file it is handed, in that file's own order,
at that file's own measure. The three `design-*` builders are alternates inherited from #77: they
arrange semantic slots ("headline", "causal") that predate the grid, so each can only show a
subset of the board and each chooses that subset itself. Useful to look at, not what gets printed.

CONTENT COMES FROM THE BOARD FILES, both of them, as of 17 August. This file arrived with both
posters' words in two hand-written dicts. #78 wrote `board_content.py` to replace them and left
them in place -- `bottom_board` was imported and never called -- so three docstrings claimed a
poster could not say something the build sheet does not, while the code still said whatever the
dicts said. The evidence was in the file it described: the factory dict's sources line cited
`reports/09-the-poster.md` and `12-poster-build-sheet-v2.md` weeks after both were renamed, and
no link checker could see it because a filename inside a Python string is not a link. The footer
citation is now built from the `Path` that was actually read, so that class of drift is not
available any more.

THE SIZE GUARD. The docstring used to say "3x4-foot" and the eight committed outputs were
36 x 48 inches; the UW vertical template is 24 x 36, 2.25x smaller in area. That was my error
before it was anybody else's -- the build sheet asserted 3ft x 4ft for a fortnight and this file
inherited it. `assert_template_size()` fails loudly rather than rendering a poster nobody can
print, because the defect is invisible until the printer refuses it.

THE "LAYOUTS ARE SIZED FOR THE LARGER BOARD" WARNING WAS OVERSTATED, and it is worth saying so
rather than deleting it. Every one of those oversized geometry blocks is in `build_evidence_grid`
and the two `landscape=True` branches -- the *horizontal* concept, whose template #78 deliberately
did not take. The three portrait builders already fit 24 in (1.50 + 21.0 = 22.50). So the fix was
not design work on seven blocks; it was dropping one builder that could not have run anyway, for
want of a file. A warning inherited without being checked costs about what a wrong constant does.
"""

from __future__ import annotations

import json
import math
import re
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import fitz
import win32com.client
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


from board_content import (
    BOTTOM,
    TOP,
    Panel,
    board,
    check as check_board,
    long_panels,
    stale_counts,
    title_block,
)

ROOT = Path(__file__).resolve().parent
TEMPLATES = ROOT / "templates"
OUTPUTS = ROOT / "outputs"
FIGURES = ROOT.parent / "reports" / "figures"

PORTRAIT_TEMPLATE = TEMPLATES / "ResearchPoster_Template_Vertical_2023.pptx"
LANDSCAPE_TEMPLATE = TEMPLATES / "ResearchPoster_Template_Horizontal_2023.pptx"

UW_PURPLE = "32006E"
UW_GOLD = "E8D3A2"
INK = "1F1F24"
MUTED = "5E5A66"
PAPER = "FFFFFF"
PANEL = "F5F3F7"
LAVENDER = "EEEAF4"
GOLD_TINT = "F8F2E4"
GREEN = "1FAE7A"
BLUE = "2F78D0"
ORANGE = "F0652F"
RED = "B63A3A"
LINE = "D9D2E3"

# The verdict colors for the speedups table's kept? column. Dark-medium on Jeffrey's spec,
# so they read as ink with a hue rather than as highlights; the words stay beside them, so no
# verdict is encoded in color alone.
KEPT_GREEN = "1F7A4C"
KEPT_RED = "B03A2E"
KEPT_BLUE = "2A5DB0"

HEADLINE_FONT = "Encode Sans Normal"
# What the template's theme actually names as its major font. We had been setting the Normal
# weight and turning bold on, which is a different, lighter letterform than the one the top board
# sets its title in -- side by side on a wall the two titles did not look like one poster.
TITLE_BLACK = "Encode Sans Normal Black"
SUBHEAD_FONT = "Uni Sans Book"
BODY_FONT = "Open Sans"


# ---------------------------------------------------------------------------------------------
# Content. Both boards' words come out of their build sheets; nothing that prints is typed here.
# ---------------------------------------------------------------------------------------------
#
# PR #77 carried both posters' words in two hand-written dicts. #78 wrote `board_content.py` to
# replace them and left them in place -- the parser was imported and never called, so three
# docstrings claimed a poster could not say something the build sheet does not, and the code did
# not do it. The evidence was sitting in the file it described: the factory dict's sources line
# cited `reports/09-the-poster.md` and `12-poster-build-sheet-v2.md` weeks after both were
# renamed, and no link checker could see it, because a filename inside a Python string is not a
# link. That is this project's own pattern one level up -- not a constant deciding a result, but
# a claim about the code that the code does not implement.

# Which cell fills each slot in the three alternate designs. Those designs predate the measured
# grid and name their slots semantically ("headline", "causal"), so the mapping is written down
# rather than guessed at. `build_board` ignores all of it and sets the cells in the order the
# board file gives them, which is why that is the one to print.
SLOTS = {
    "yoruba": {
        "question": ("1",),
        "data": ("2",),
        "tokenizer": ("3", "4"),
        "headline": ("5",),
        "ner": ("6", "7"),
        "causal": ("8",),
        "variance": ("9",),
        "limits": ("strip2",),
    },
}

# The four big numbers a design sets in its stat row, as cell numbers. The factory board has no
# entry: v3 replaced its nine big-number pills with one stat rail, whose six rows come off panel
# 2's own body rather than out of a tuple here.
STAT_CELLS = {"yoruba": ("5", "3", "1", "9")}

EYEBROW = {
    "yoruba": "LOW-RESOURCE LANGUAGE MODELING",
    "factory": "EXPERIMENT INFRASTRUCTURE",
}

BOARD_FILE = {"yoruba": TOP, "factory": BOTTOM}


def content_from_board(kind: str) -> dict:
    """Everything a builder needs, read from the board file for that half of the poster."""
    path = BOARD_FILE[kind]
    panels = board(path)
    head = title_block(path)

    def joined(*keys: str) -> str:
        return "\n\n".join(panels[k].body_plain for k in keys if k in panels)

    content = {
        "title": head.title,
        "subtitle": head.subtitle,
        "eyebrow": EYEBROW[kind],
        "takeaway": head.takeaway,
        "author": head.author,
        "goals": head.goals,
        "sources": head.sources_line,
        "stats": [
            (panels[c].big_lines[0], panels[c].big_lines[-1])
            for c in STAT_CELLS.get(kind, ())
            if c in panels and panels[c].big_lines
        ],
        "panels": panels,
    }
    for slot, cells in SLOTS.get(kind, {}).items():
        content[slot] = joined(*cells)
    return content


# ---------------------------------------------------------------------------------------------
# The grid, measured off the template in report 12 rather than chosen here.
# ---------------------------------------------------------------------------------------------
COL_X = (1.50, 8.75, 16.13)
COL_W = 6.35
ROW_Y = (9.25, 16.20, 23.15)
ROW_H = 6.70
STRIP_Y = 30.10
STRIP_H = 4.45

# What a cell spends its 6.70 in on. Report 12 budgets ~0.6 header, ~1.4 big number, ~2.3 figure
# and ~1.9 body; these are that plan adjusted to what PowerPoint actually laid out, measured off
# `layout-validation.json` rather than estimated.
#
# THE FIGURE BLOCK IS 2.30 in, NOT REPORT 12's 2.30-BY-GUESS: it is what the column figures are
# drawn at, so they land at scale 1.0 and the type inside them prints at the size it was set.
# Scaled to the old 1.80 in they came back to 0.75x and 7 pt. The big number gives up 0.15 in for
# it, being decoration next to a chart.
#
# THE MEASURED RATE IS ~2.82 pt OF COLUMN PER WORD at 18 pt in a 5.77 in measure. So 1.9 in of
# body is 137 pt, which holds about **48 words, not 55** -- report 12's guide overshoots its own
# geometry by roughly 13%, and every cell on both boards was written to the guide. Rather than
# re-cut two boards to 48 words, the body keeps 2.72 in (196 pt, ~69 words), taken out of the
# figure and the big number. Both are decorative next to the prose.
CELL_TOP = 0.12
CELL_HEAD_H = 0.70
CELL_BIG_H = 0.80
CELL_FIG_H = 2.12
CELL_BOTTOM = 0.15
GAP = 0.06

# The strip is 4.45 in rather than 6.70 and cannot hold 18 pt prose: a strip block carrying a
# figure would have 1.50 in of body left, which is 38 words against the ~100 report 12 budgets.
# The strip is the foot of the board and is read last, so it is set two steps down instead.
STRIP_FIG_H = 1.20
STRIP_BODY_PT = 15.0

BODY_PT = 18.0
HEAD_PT = 20.0


@dataclass(frozen=True)
class FigureSpec:
    path: Path
    x: float
    y: float
    w: float
    h: float
    name: str


def assert_template_size(prs, expect=(24.0, 36.0), name="template") -> None:
    """Refuse to build a poster at a size nobody can print.

    PR #77 shipped eight posters at 36 x 48 because the build sheet said 3ft x 4ft and nobody
    measured the template. The failure is silent all the way to the print shop, so it gets an
    assertion rather than a comment: a wrong-size poster looks exactly like a right-size one on
    screen, which is the whole problem.
    """
    w = prs.slide_width / 914400
    h = prs.slide_height / 914400
    if (round(w, 2), round(h, 2)) != expect:
        raise SystemExit(
            f"{name} is {w:.2f} x {h:.2f} in; expected {expect[0]} x {expect[1]}. "
            f"The UW vertical template is 24 x 36. If this is deliberate, change `expect` "
            f"and re-fit the layouts -- they place content out to x = 34.75 in, which does "
            f"not fit a 24 in board."
        )


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[index]


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def open_template(path: Path, slide_index: int = 0) -> tuple[Presentation, object]:
    prs = Presentation(path)
    for index in reversed(range(len(prs.slides))):
        if index != slide_index:
            delete_slide(prs, index)
    slide = prs.slides[0]
    clear_slide(slide)
    return prs, slide


def set_cell_margins(text_frame, margin: float = 0.08) -> None:
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float,
    *,
    color: str = INK,
    font: str = BODY_FONT,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.03,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    # python-pptx writes <a:spAutoFit/> into every textbox it creates, which is
    # MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT: the shape grows downward until the text fits. That makes
    # the height passed in above a suggestion rather than a constraint, and it silently disabled
    # the one gate this module has. `embed_figures_and_export` flags a box whose text is taller
    # than the box -- but with autofit on, PowerPoint had already grown the box to match, so
    # BoundHeight was always a few points UNDER Height and the check reported clean on a board
    # whose cells visibly ran over each other. 55 shapes examined, 0 skipped, 0 flagged.
    # Turning autofit off makes the declared geometry authoritative and the overflow check real.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    set_cell_margins(tf, margin)
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = PAPER,
    line: str | None = None,
    radius: bool = False,
    name: str | None = None,
):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    if name:
        shape.name = name
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color=UW_GOLD, width=1.2):
    line = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_header(
    slide,
    content: dict,
    *,
    landscape: bool = False,
    subtitle_pt: float = 18.0,
    author_pt: float = 9.5,
    title_pt: float = 48.0,
    title_font: str | None = None,
    title_caps: bool = False,
) -> None:
    """The purple band: eyebrow or wordmark, title, subtitle, author line.

    The two type sizes are parameters with their inherited values as defaults. The bottom board
    raises both -- 9.5 pt on a 2 x 3 ft board is not a small line, it is an absent one, and that
    line carries the team names the assignment asks for. The top board is not mine to restyle, so
    it keeps what it had by taking the defaults.
    """
    if landscape:
        add_text(
            slide,
            1.25,
            0.82,
            4.7,
            0.35,
            content["eyebrow"],
            10,
            color=UW_GOLD,
            font=SUBHEAD_FONT,
            bold=True,
            name="Header_Eyebrow",
        )
        add_text(
            slide,
            1.25,
            1.12,
            28.4,
            1.62,
            content["title"],
            50,
            color=PAPER,
            font=HEADLINE_FONT,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            name="Header_Title",
        )
        add_text(
            slide,
            1.28,
            2.78,
            26.5,
            0.55,
            content["subtitle"],
            20,
            color=PAPER,
            font=SUBHEAD_FONT,
            name="Header_Subtitle",
        )
    else:
        # The vertical template master carries the UW wordmark and a tiny
        # "research poster template" label above it. Cover only that label.
        add_rect(slide, 1.42, 1.10, 4.35, 0.30, fill=UW_PURPLE, name="MasterLabelCover")
        add_text(
            slide,
            1.50,
            1.80,
            19.3,
            2.50,
            content["title"].upper() if title_caps else content["title"],
            title_pt,
            color=PAPER,
            font=title_font or HEADLINE_FONT,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
            name="Header_Title",
        )
        add_text(
            slide,
            1.53,
            4.45,
            19.0,
            0.95,
            content["subtitle"],
            subtitle_pt,
            color=PAPER,
            font=SUBHEAD_FONT,
            name="Header_Subtitle",
        )
        add_text(
            slide,
            1.53,
            7.12,
            19.0,
            0.60,
            # The rubric asks for team member names, so this comes off the board's author line
            # rather than being a generic string. It said "A2-NLP Project Team" for a fortnight.
            # Same FACE as the subtitle, not just the same size -- 19 pt of Open Sans against
            # 19 pt of Uni Sans read as two different sizes, which was Jeffrey's exact note.
            content.get("author")
            or "CSED 504 · University of Washington · Summer 2026",
            author_pt,
            valign=MSO_ANCHOR.MIDDLE,
            color=UW_GOLD,
            font=SUBHEAD_FONT,
            name="Header_Author",
        )


def add_takeaway(slide, content: dict, x: float, y: float, w: float, h: float,
                 *, label: str = "TAKEAWAY") -> None:
    add_rect(slide, x, y, w, h, fill=GOLD_TINT, line=UW_GOLD, name="Takeaway_Background")
    add_text(
        slide,
        x + 0.28,
        y + 0.16,
        1.55,
        h - 0.28,
        label,
        10.5,
        color=UW_PURPLE,
        font=SUBHEAD_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        name="Takeaway_Label",
    )
    add_text(
        slide,
        x + 1.72,
        y + 0.15,
        w - 2.0,
        h - 0.25,
        content["takeaway"],
        15.5,
        color=INK,
        font=BODY_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        name="Takeaway_Text",
    )


def add_panel(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    heading: str,
    body: str,
    *,
    fill: str = PANEL,
    accent: str = UW_PURPLE,
    body_size: float = 13.3,
    heading_size: float = 19.0,
    label: str | None = None,
    name: str = "Panel",
):
    add_rect(slide, x, y, w, h, fill=fill, line=LINE, radius=True, name=f"{name}_Box")
    add_rect(slide, x, y, 0.11, h, fill=accent, name=f"{name}_Accent")
    if label:
        add_text(
            slide,
            x + 0.30,
            y + 0.23,
            w - 0.58,
            0.28,
            label.upper(),
            8.7,
            color=accent,
            font=SUBHEAD_FONT,
            bold=True,
            name=f"{name}_Label",
        )
        heading_y = y + 0.60
    else:
        heading_y = y + 0.27
    add_text(
        slide,
        x + 0.30,
        heading_y,
        w - 0.58,
        0.68,
        heading,
        heading_size,
        color=accent,
        font=HEADLINE_FONT,
        bold=True,
        name=f"{name}_Heading",
    )
    add_text(
        slide,
        x + 0.30,
        heading_y + 0.83,
        w - 0.58,
        h - (heading_y - y) - 1.03,
        body,
        body_size,
        color=INK,
        font=BODY_FONT,
        name=f"{name}_Body",
    )


def add_stat(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    *,
    fill: str = UW_PURPLE,
    value_color: str = PAPER,
    label_color: str = UW_GOLD,
    value_size: float = 29,
    name: str = "Stat",
):
    add_rect(slide, x, y, w, h, fill=fill, radius=True, name=f"{name}_Box")
    # The value box and the label box must not share vertical space. They used to -- value from
    # y+0.13 to y+0.13+h*0.55 and label from y+h*0.60, an overlap of ~0.09 in on the 0.80 in cell
    # band -- and it went unseen for a whole review cycle because every big number on both boards
    # was digits, and digits have no descenders. "fingerprint" does, and its g and p ran straight
    # through the gold hash line under it. The overflow gate cannot catch this class: it measures
    # each shape against its own box, and both shapes fit; the collision is BETWEEN shapes.
    # ...and the value band is sized from the TYPE IT WAS GIVEN, not from a fraction of the cell.
    # h*0.52 is 30.0 pt on the 0.80 in cell band, while one line of the 27 pt value needs 32.6 pt
    # under the font PowerPoint substitutes on this machine. The fraction encodes an assumption
    # about font metrics, and the fonts are deliberately not in the repo (#78 dropped them as a
    # commercial licence), so every renderer gets whatever its own PowerPoint picks. Every big
    # number on both boards overflowed by 2.6 pt here while rendering clean on the workstation --
    # same code, same board, different substitute. Deriving the height from value_size makes the
    # band correct wherever it is rendered; the label sits under whatever that came to, which
    # keeps the separation #106 established rather than re-deriving it from h.
    value_h = max(h * 0.52, value_size * 1.28 / 72)
    label_y = y + 0.08 + value_h + 0.02
    label_h = max(0.14, (y + h) - label_y - 0.05)
    add_text(
        slide,
        x + 0.14,
        y + 0.08,
        w - 0.28,
        value_h,
        value,
        value_size,
        color=value_color,
        font=HEADLINE_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        name=f"{name}_Value",
    )
    add_text(
        slide,
        x + 0.18,
        label_y,
        w - 0.36,
        label_h,
        label,
        9.4,
        color=label_color,
        font=BODY_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        name=f"{name}_Label",
    )


def add_figure_frame(
    slide,
    figure_specs: list[FigureSpec],
    filename: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str | None = None,
    caption: str | None = None,
    fill: str = PAPER,
    pad: float = 0.15,
    name: str = "Figure",
) -> None:
    add_rect(slide, x, y, w, h, fill=fill, line=LINE, radius=True, name=f"{name}_Frame")
    label_h = 0.42 if label else 0
    caption_h = 0.46 if caption else 0
    if label:
        add_text(
            slide,
            x + 0.22,
            y + 0.12,
            w - 0.44,
            0.27,
            label.upper(),
            8.6,
            color=UW_PURPLE,
            font=SUBHEAD_FONT,
            bold=True,
            name=f"{name}_Label",
        )
    if caption:
        add_text(
            slide,
            x + 0.24,
            y + h - 0.42,
            w - 0.48,
            0.28,
            caption,
            8.7,
            color=MUTED,
            font=BODY_FONT,
            italic=True,
            align=PP_ALIGN.CENTER,
            name=f"{name}_Caption",
        )
    ext = Path(filename).suffix
    source = FIGURES / filename
    if not source.exists() and ext.lower() == ".svg":
        source = FIGURES / f"{Path(filename).stem}.png"
    figure_specs.append(
        FigureSpec(
            source,
            x + pad,
            y + pad + label_h,
            w - 2 * pad,
            h - 2 * pad - label_h - caption_h,
            name,
        )
    )


def add_footer(slide, content: dict, *, landscape: bool = False,
               footer_y: float | None = None) -> None:
    if landscape:
        y = 23.36
        add_line(slide, 1.25, y, 34.75, y, UW_GOLD, 1)
        add_text(
            slide,
            1.25,
            y + 0.10,
            31.0,
            0.34,
            content["sources"],
            7.4,
            color=MUTED,
            font=BODY_FONT,
            name="Footer_Sources",
        )
        add_text(
            slide,
            32.2,
            y + 0.10,
            2.55,
            0.34,
            "CSED 504 · 2026",
            7.8,
            color=UW_PURPLE,
            font=SUBHEAD_FONT,
            bold=True,
            align=PP_ALIGN.RIGHT,
            name="Footer_Tag",
        )
    else:
        # 34.84 and 12 pt, not 35.16 and 7.2. The footer carries the AI statement now, and a
        # compliance line set at 7.2 pt on a board two feet wide is not a compliance line.
        y = footer_y if footer_y is not None else 34.84
        add_line(slide, 1.50, y, 22.50, y, UW_GOLD, 1)
        add_text(
            slide,
            1.50,
            y + 0.12,
            18.3,
            1.30,
            content["sources"],
            FACTORY_PT["footer"],
            color=MUTED,
            font=BODY_FONT,
            name="Footer_Sources",
        )
        # Inside the footer band, not above the rule: above it the tag sat in row C's last
        # column and printed through the final source reference on two proofs running.
        add_text(
            slide,
            20.0,
            y + 0.14,
            2.5,
            0.30,
            "CSED 504 · 2026",
            10.0,
            color=UW_PURPLE,
            font=SUBHEAD_FONT,
            bold=True,
            align=PP_ALIGN.RIGHT,
            name="Footer_Tag",
        )


def build_classic(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    prs, slide = open_template(PORTRAIT_TEMPLATE, 0)
    figures: list[FigureSpec] = []
    add_header(slide, content)
    add_takeaway(slide, content, 1.5, 7.65, 21.0, 1.16)

    if kind == "yoruba":
        add_panel(
            slide,
            1.5,
            9.20,
            6.5,
            8.05,
            "The choice",
            content["question"] + "\n\n" + content["headline"],
            label="Research question",
            body_size=15.5,
            name="Choice",
        )
        add_figure_frame(
            slide,
            figures,
            "01-headline.svg",
            8.4,
            9.20,
            14.1,
            8.05,
            label="Headline result",
            caption="SIB-200 uses dev-selected rates; NER does not.",
            name="HeadlineFigure",
        )
        add_figure_frame(
            slide,
            figures,
            "02-tokenizer-gradient.svg",
            1.5,
            17.72,
            10.28,
            7.65,
            label="Coverage gradient",
            caption="Tokenizer fertility across 17 languages.",
            name="GradientFigure",
        )
        add_figure_frame(
            slide,
            figures,
            "05-data-saturation.svg",
            12.22,
            17.72,
            10.28,
            7.65,
            label="Data control",
            caption="More text stops moving loss at these budgets.",
            name="DataFigure",
        )
        add_figure_frame(
            slide,
            figures,
            "17-tokenizer-lottery.svg",
            1.5,
            25.84,
            10.28,
            8.82,
            label="Run-to-run variability",
            caption="Six registered seeds separate variance, not mean.",
            name="LotteryFigure",
        )
        add_panel(
            slide,
            12.22,
            25.84,
            10.28,
            8.82,
            "What survives",
            content["tokenizer"]
            + "\n\n"
            + content["causal"]
            + "\n\n"
            + "Boundary: "
            + content["limits"],
            fill=LAVENDER,
            body_size=14.0,
            heading_size=20,
            name="Conclusions",
        )
    else:
        add_panel(
            slide,
            1.5,
            9.20,
            6.5,
            8.05,
            "Why a factory?",
            content["question"] + "\n\n" + content["workflow"],
            label="System goal",
            body_size=15.5,
            name="FactoryQuestion",
        )
        add_figure_frame(
            slide,
            figures,
            "14-where-the-speedup-came-from.svg",
            8.4,
            9.20,
            14.1,
            8.05,
            label="Measured speedup",
            caption="Efficiency and parallelism are reported separately.",
            name="SpeedFigure",
        )
        add_figure_frame(
            slide,
            figures,
            "15-what-a-run-is-made-of.svg",
            1.5,
            17.72,
            10.28,
            7.65,
            label="Notebook / queue boundary",
            caption="Cheap stages stay interactive; training becomes a job.",
            name="PipelineFigure",
        )
        add_figure_frame(
            slide,
            figures,
            "19-the-interface.svg",
            12.22,
            17.72,
            10.28,
            7.65,
            label="Public surface",
            caption="Nine stable calls hide the machinery.",
            name="InterfaceFigure",
        )
        add_figure_frame(
            slide,
            figures,
            "07-dashboard.png",
            1.5,
            25.84,
            10.28,
            8.82,
            label="Operations",
            caption="Queue, GPU state and comparable run records.",
            name="DashboardFigure",
        )
        add_panel(
            slide,
            12.22,
            25.84,
            10.28,
            8.82,
            "What made it trustworthy",
            content["records"]
            + "\n\n"
            + content["statistics"]
            + "\n\n"
            + content["honesty"],
            fill=LAVENDER,
            body_size=14.0,
            heading_size=20,
            name="Trust",
        )
    add_footer(slide, content)
    return prs, figures


def build_spine(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    prs, slide = open_template(PORTRAIT_TEMPLATE, 1)
    figures: list[FigureSpec] = []
    add_header(slide, content)
    add_takeaway(slide, content, 1.5, 7.45, 21.0, 1.16)

    left_x, mid_x, right_x = 1.5, 8.75, 15.85
    left_w, mid_w, right_w = 6.8, 6.65, 6.65
    body_y, body_h = 9.02, 25.62
    add_rect(slide, mid_x, body_y, mid_w, body_h, fill=UW_PURPLE, name="NarrativeSpine")

    if kind == "yoruba":
        add_figure_frame(
            slide,
            figures,
            "02-tokenizer-gradient.svg",
            left_x,
            body_y,
            left_w,
            8.45,
            label="Why vocabulary fit matters",
            name="SpineGradient",
        )
        add_panel(
            slide,
            left_x,
            17.85,
            left_w,
            6.25,
            "Enough text",
            content["data"],
            body_size=13.5,
            heading_size=18,
            name="SpineData",
        )
        add_figure_frame(
            slide,
            figures,
            "05-data-saturation.svg",
            left_x,
            24.48,
            left_w,
            10.16,
            label="Fixed-compute control",
            name="SpineDataFigure",
        )
        add_text(
            slide,
            mid_x + 0.40,
            body_y + 0.50,
            mid_w - 0.8,
            0.46,
            "THE EVIDENCE CHAIN",
            10,
            color=UW_GOLD,
            font=SUBHEAD_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="SpineLabel",
        )
        for i, (value, label) in enumerate(content["stats"]):
            add_text(
                slide,
                mid_x + 0.45,
                body_y + 1.35 + i * 3.55,
                mid_w - 0.9,
                1.25,
                value,
                31,
                color=PAPER,
                font=HEADLINE_FONT,
                bold=True,
                align=PP_ALIGN.CENTER,
                valign=MSO_ANCHOR.MIDDLE,
                name=f"SpineStat{i}_Value",
            )
            add_text(
                slide,
                mid_x + 0.6,
                body_y + 2.58 + i * 3.55,
                mid_w - 1.2,
                0.67,
                label,
                10.2,
                color=UW_GOLD,
                font=BODY_FONT,
                bold=True,
                align=PP_ALIGN.CENTER,
                name=f"SpineStat{i}_Label",
            )
            if i < 3:
                add_line(
                    slide,
                    mid_x + 1.5,
                    body_y + 4.47 + i * 3.55,
                    mid_x + mid_w - 1.5,
                    body_y + 4.47 + i * 3.55,
                    UW_GOLD,
                    0.7,
                )
        add_text(
            slide,
            mid_x + 0.48,
            body_y + 16.25,
            mid_w - 0.96,
            7.8,
            "MEAN\nNo reliable matched-compute difference in pretraining bits per character "
            "(Welch p=0.374).\n\n"
            "VARIANCE\nThe 250k-vocabulary arm was nearly four times as variable "
            "(sd 0.145 vs 0.037; p=0.0098).\n\n"
            "DOWNSTREAM\nHolding compute fixed, the fitted vocabulary led by +0.144 on SIB-200 "
            "and +0.061 on NER; every seed won.\n\n"
            "NEXT TEST\nRepeat the matched-compute vocabulary swap across the coverage gradient. "
            "Causal downstream evidence is still Yoruba-only.",
            13.8,
            color=PAPER,
            font=BODY_FONT,
            name="SpineInterpretation",
        )
        add_figure_frame(
            slide,
            figures,
            "01-headline.svg",
            right_x,
            body_y,
            right_w,
            9.28,
            label="Downstream result",
            name="SpineHeadline",
        )
        add_panel(
            slide,
            right_x,
            18.70,
            right_w,
            6.22,
            "Tasks disagree",
            content["ner"],
            body_size=13.2,
            heading_size=18,
            name="SpineTasks",
        )
        add_figure_frame(
            slide,
            figures,
            "17-tokenizer-lottery.svg",
            right_x,
            25.30,
            right_w,
            9.34,
            label="Variance, not mean",
            name="SpineLottery",
        )
    else:
        add_figure_frame(
            slide,
            figures,
            "14-where-the-speedup-came-from.svg",
            left_x,
            body_y,
            left_w,
            8.45,
            label="Measure the speedup",
            name="SpineSpeed",
        )
        add_panel(
            slide,
            left_x,
            17.85,
            left_w,
            6.25,
            "Split by cost",
            content["pipeline"],
            body_size=13.5,
            heading_size=18,
            name="SpinePipeline",
        )
        add_figure_frame(
            slide,
            figures,
            "15-what-a-run-is-made-of.svg",
            left_x,
            24.48,
            left_w,
            10.16,
            label="Interactive → unattended",
            name="SpinePipelineFigure",
        )
        add_text(
            slide,
            mid_x + 0.4,
            body_y + 0.50,
            mid_w - 0.8,
            0.46,
            "THE FACTORY LOOP",
            10,
            color=UW_GOLD,
            font=SUBHEAD_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="SpineLabel",
        )
        steps = [
            ("01", "PREPARE", "tokenize once; fingerprint"),
            ("02", "ESTIMATE", "measure this GPU"),
            ("03", "PRETRAIN", "one cell or a fleet"),
            ("04", "READ", "curves and result records"),
        ]
        for i, (num, title, desc) in enumerate(steps):
            sy = body_y + 1.45 + i * 3.55
            add_text(
                slide,
                mid_x + 0.55,
                sy,
                1.15,
                0.72,
                num,
                16,
                color=UW_GOLD,
                font=HEADLINE_FONT,
                bold=True,
                name=f"Step{i}_Num",
            )
            add_text(
                slide,
                mid_x + 1.65,
                sy,
                mid_w - 2.1,
                0.58,
                title,
                16,
                color=PAPER,
                font=HEADLINE_FONT,
                bold=True,
                name=f"Step{i}_Title",
            )
            add_text(
                slide,
                mid_x + 1.65,
                sy + 0.72,
                mid_w - 2.1,
                0.62,
                desc,
                9.7,
                color=UW_GOLD,
                font=BODY_FONT,
                name=f"Step{i}_Desc",
            )
            if i < 3:
                add_line(
                    slide,
                    mid_x + mid_w / 2,
                    sy + 1.73,
                    mid_x + mid_w / 2,
                    sy + 3.1,
                    UW_GOLD,
                    1,
                )
        add_text(
            slide,
            mid_x + 0.5,
            body_y + 16.25,
            mid_w - 1.0,
            7.8,
            "THE RESULT\nThe claims gate reports 6 supported, 2 unsupported and 1 underpowered "
            "claim. Failures stay visible.\n\n"
            "THE STANDARD\nA result must be reproducible enough for another person to challenge "
            "and overturn it. That happened twice—and is evidence the tooling worked.",
            15.0,
            color=PAPER,
            font=BODY_FONT,
            name="SpineFactoryResult",
        )
        add_figure_frame(
            slide,
            figures,
            "19-the-interface.svg",
            right_x,
            body_y,
            right_w,
            9.28,
            label="Nine-function surface",
            name="SpineInterface",
        )
        add_panel(
            slide,
            right_x,
            18.70,
            right_w,
            6.22,
            "Records survive",
            content["records"],
            body_size=13.2,
            heading_size=18,
            name="SpineRecords",
        )
        add_figure_frame(
            slide,
            figures,
            "07-dashboard.png",
            right_x,
            25.30,
            right_w,
            9.34,
            label="One operational view",
            name="SpineDashboard",
        )
    add_footer(slide, content)
    return prs, figures


def build_evidence_grid(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    """UNREACHABLE, and kept deliberately rather than deleted. Not in BUILDERS.

    This is the landscape concept. It opens `LANDSCAPE_TEMPLATE`, which #78 deliberately did not
    take from PR #77 on the grounds that neither board is landscape -- so the file is not in the
    repository and this function raises before it draws anything. It is also the source of every
    "content runs out to x = 34.75 in" warning in this module: those blocks are correct for a
    48 in board and were never a defect in the portrait path.

    Kept because it is the only worked-out wide layout in the project, and restoring it needs one
    template file rather than a rewrite. Delete it if the horizontal template is ruled out for
    good.
    """
    prs, slide = open_template(LANDSCAPE_TEMPLATE, 0)
    figures: list[FigureSpec] = []
    add_header(slide, content, landscape=True)
    add_takeaway(slide, content, 1.25, 4.25, 33.5, 1.05)

    if kind == "yoruba":
        add_figure_frame(
            slide,
            figures,
            "01-headline.svg",
            1.25,
            5.72,
            19.9,
            9.12,
            label="Result",
            name="GridHeadline",
        )
        add_panel(
            slide,
            21.58,
            5.72,
            13.17,
            5.20,
            "A small model, a precise claim",
            content["headline"],
            fill=LAVENDER,
            body_size=16.0,
            heading_size=20,
            name="GridClaim",
        )
        add_panel(
            slide,
            21.58,
            11.34,
            13.17,
            3.50,
            "Causal vocabulary swap",
            content["causal"],
            body_size=13.5,
            heading_size=17,
            name="GridTask",
        )
        add_figure_frame(
            slide,
            figures,
            "02-tokenizer-gradient.svg",
            1.25,
            15.28,
            10.78,
            7.55,
            label="1 · Coverage",
            name="GridGradient",
        )
        add_figure_frame(
            slide,
            figures,
            "05-data-saturation.svg",
            12.45,
            15.28,
            10.78,
            7.55,
            label="2 · Data",
            name="GridData",
        )
        add_figure_frame(
            slide,
            figures,
            "17-tokenizer-lottery.svg",
            23.65,
            15.28,
            11.10,
            7.55,
            label="3 · Variability",
            name="GridLottery",
        )
    else:
        add_figure_frame(
            slide,
            figures,
            "14-where-the-speedup-came-from.svg",
            1.25,
            5.72,
            17.0,
            9.12,
            label="Measure",
            name="GridSpeed",
        )
        add_panel(
            slide,
            18.68,
            5.72,
            16.07,
            4.48,
            "From notebook state to evidence",
            content["question"] + "\n\n" + content["speed"],
            fill=LAVENDER,
            body_size=14.5,
            heading_size=20,
            name="GridFactoryIntro",
        )
        stat_w = 3.68
        for i, (value, label) in enumerate(content["stats"]):
            add_stat(
                slide,
                18.68 + i * (stat_w + 0.43),
                10.65,
                stat_w,
                4.19,
                value,
                label,
                value_size=24,
                name=f"GridStat{i}",
            )
        add_figure_frame(
            slide,
            figures,
            "15-what-a-run-is-made-of.svg",
            1.25,
            15.28,
            10.78,
            7.55,
            label="1 · Split by cost",
            name="GridPipeline",
        )
        add_figure_frame(
            slide,
            figures,
            "19-the-interface.svg",
            12.45,
            15.28,
            10.78,
            7.55,
            label="2 · Hide complexity",
            name="GridInterface",
        )
        add_figure_frame(
            slide,
            figures,
            "13-how-many-seeds.svg",
            23.65,
            15.28,
            11.10,
            7.55,
            label="3 · Limit the claim",
            name="GridSeeds",
        )
    add_footer(slide, content, landscape=True)
    return prs, figures


def build_editorial(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    prs, slide = open_template(PORTRAIT_TEMPLATE, 0)
    figures: list[FigureSpec] = []
    add_header(slide, content)
    add_takeaway(slide, content, 1.5, 7.65, 21.0, 1.16)

    if kind == "yoruba":
        add_figure_frame(
            slide,
            figures,
            "01-headline.svg",
            1.5,
            9.20,
            14.4,
            10.35,
            label="The result",
            name="EditorialHeadline",
        )
        add_rect(slide, 16.32, 9.20, 6.18, 10.35, fill=UW_PURPLE, name="EditorialStatsBox")
        add_text(
            slide,
            16.78,
            9.72,
            5.25,
            0.45,
            "READ THE RESULT AS A PAIR",
            9.2,
            color=UW_GOLD,
            font=SUBHEAD_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="EditorialStatsLabel",
        )
        add_text(
            slide,
            16.82,
            10.55,
            5.15,
            1.40,
            "0.688",
            34,
            color=PAPER,
            font=HEADLINE_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="EditorialStat1",
        )
        add_text(
            slide,
            16.82,
            11.88,
            5.15,
            1.05,
            "topic macro-F1\nfrom scratch",
            10.3,
            color=UW_GOLD,
            font=BODY_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="EditorialStat1Label",
        )
        add_line(slide, 17.5, 13.25, 21.3, 13.25, UW_GOLD, 0.8)
        add_text(
            slide,
            16.82,
            13.62,
            5.15,
            1.40,
            "0.837",
            34,
            color=PAPER,
            font=HEADLINE_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="EditorialStat2",
        )
        add_text(
            slide,
            16.82,
            14.95,
            5.15,
            1.05,
            "entity F1\nfrom scratch",
            10.3,
            color=UW_GOLD,
            font=BODY_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            name="EditorialStat2Label",
        )
        add_text(
            slide,
            16.75,
            16.55,
            5.35,
            2.25,
            "Ahead of mmBERT on topic; slightly behind on entities. The disagreement is part of the finding.",
            11.0,
            color=PAPER,
            font=BODY_FONT,
            align=PP_ALIGN.CENTER,
            name="EditorialStatNote",
        )
        add_figure_frame(
            slide,
            figures,
            "02-tokenizer-gradient.svg",
            1.5,
            20.05,
            13.15,
            8.40,
            label="The mechanism",
            name="EditorialGradient",
        )
        add_panel(
            slide,
            15.08,
            20.05,
            7.42,
            8.40,
            "Why vocabulary fit",
            content["tokenizer"] + "\n\n" + content["causal"],
            fill=GOLD_TINT,
            accent=ORANGE,
            body_size=13.2,
            heading_size=19,
            name="EditorialMechanism",
        )
        col_w = 6.72
        add_panel(
            slide,
            1.5,
            28.90,
            col_w,
            5.76,
            "Text is not the bottleneck",
            content["data"],
            body_size=13.5,
            heading_size=16.5,
            name="EditorialData",
        )
        add_panel(
            slide,
            8.64,
            28.90,
            col_w,
            5.76,
            "Means and variability",
            content["variance"],
            body_size=13.2,
            heading_size=16.5,
            name="EditorialVariance",
        )
        add_panel(
            slide,
            15.78,
            28.90,
            col_w,
            5.76,
            "Scope",
            content["limits"],
            body_size=13.2,
            heading_size=16.5,
            name="EditorialScope",
        )
    else:
        add_figure_frame(
            slide,
            figures,
            "19-the-interface.svg",
            1.5,
            9.20,
            14.4,
            10.35,
            label="One screen to operate the factory",
            name="EditorialInterface",
        )
        add_rect(slide, 16.32, 9.20, 6.18, 10.35, fill=UW_PURPLE, name="EditorialStatsBox")
        for i, (value, label) in enumerate(content["stats"][:3]):
            sy = 9.62 + i * 3.18
            add_text(
                slide,
                16.72,
                sy,
                5.38,
                1.2,
                value,
                31,
                color=PAPER,
                font=HEADLINE_FONT,
                bold=True,
                align=PP_ALIGN.CENTER,
                name=f"EditorialFactoryStat{i}",
            )
            add_text(
                slide,
                16.72,
                sy + 1.2,
                5.38,
                0.82,
                label,
                9.6,
                color=UW_GOLD,
                font=BODY_FONT,
                bold=True,
                align=PP_ALIGN.CENTER,
                name=f"EditorialFactoryStat{i}Label",
            )
            if i < 2:
                add_line(slide, 17.5, sy + 2.55, 21.3, sy + 2.55, UW_GOLD, 0.8)
        add_figure_frame(
            slide,
            figures,
            "14-where-the-speedup-came-from.svg",
            1.5,
            20.05,
            13.15,
            8.40,
            label="Iteration became affordable",
            name="EditorialSpeed",
        )
        add_panel(
            slide,
            15.08,
            20.05,
            7.42,
            8.40,
            "The operating rule",
            content["workflow"] + "\n\n" + content["pipeline"],
            fill=GOLD_TINT,
            accent=ORANGE,
            body_size=13.2,
            heading_size=19,
            name="EditorialWorkflow",
        )
        col_w = 6.72
        add_panel(
            slide,
            1.5,
            28.90,
            col_w,
            5.76,
            "Identity",
            content["records"],
            body_size=13.2,
            heading_size=16.5,
            name="EditorialRecords",
        )
        add_panel(
            slide,
            8.64,
            28.90,
            col_w,
            5.76,
            "Statistical guardrails",
            content["statistics"],
            body_size=13.0,
            heading_size=16.5,
            name="EditorialStats",
        )
        add_panel(
            slide,
            15.78,
            28.90,
            col_w,
            5.76,
            "Failure is evidence",
            "Healthy and doomed runs overlapped at all 11 early checkpoints, so the proposed "
            "detector could not work. The claims gate reports 6 supported, 2 unsupported and "
            "1 underpowered claim—and keeps the failures visible.",
            body_size=13.0,
            heading_size=16.5,
            name="EditorialFailure",
        )
    add_footer(slide, content)
    return prs, figures


def set_cell(
    slide,
    figures: list[FigureSpec],
    panel: Panel,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    name: str,
) -> None:
    """One cell of the board: heading, big number, figure if it has one, then the panel text.

    The vertical order is report 12's, and so are the heights. Nothing here chooses what a cell
    says or how much room it gets -- both come out of the build sheet.
    """
    strip = panel.strip
    add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True, name=f"{name}_Box")
    add_rect(slide, x, y, 0.11, h, fill=UW_PURPLE, name=f"{name}_Accent")

    inner_x, inner_w = x + 0.30, w - 0.58
    cur = y + CELL_TOP
    add_text(
        slide,
        inner_x,
        cur,
        inner_w,
        CELL_HEAD_H,
        panel.title,
        HEAD_PT,
        color=UW_PURPLE,
        font=HEADLINE_FONT,
        bold=True,
        name=f"{name}_Heading",
    )
    cur += CELL_HEAD_H + GAP

    if panel.big_lines:
        add_stat(
            slide,
            inner_x,
            cur,
            inner_w,
            CELL_BIG_H,
            panel.big_lines[0],
            panel.big_lines[-1] if len(panel.big_lines) > 1 else "",
            value_size=27,
            name=f"{name}_Big",
        )
        cur += CELL_BIG_H + GAP

    if panel.figure:
        fig_h = STRIP_FIG_H if strip else CELL_FIG_H
        add_figure_frame(
            slide,
            figures,
            panel.figure,
            inner_x,
            cur,
            inner_w,
            fig_h,
            fill=PAPER,
            pad=0.09,
            name=f"{name}_Figure",
        )
        cur += fig_h + GAP

    add_text(
        slide,
        inner_x,
        cur,
        inner_w,
        (y + h - CELL_BOTTOM) - cur,
        panel.body_plain,
        STRIP_BODY_PT if strip else BODY_PT,
        color=INK,
        font=BODY_FONT,
        name=f"{name}_Body",
    )


def build_board(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    """The board as report 12 measures it: nine cells in a 3x3, three strip blocks at the foot.

    THIS IS THE ONE TO PRINT. The three designs below it are alternates inherited from PR #77;
    they arrange semantic slots ("headline", "causal") that predate the grid, so they can only
    ever show a subset of the board and they choose that subset themselves. This one sets every
    cell the build sheet defines, in the build sheet's order, at the build sheet's measure --
    which means the printed poster and the board file cannot disagree.
    """
    prs, slide = open_template(PORTRAIT_TEMPLATE, 0)
    figures: list[FigureSpec] = []
    add_header(slide, content)
    if content.get("goals"):
        # The rubric's "Goals" item, set in the header band under the author line. It lives on
        # the board file as its own blockquote, which is why the parser keeps quote runs separate
        # -- merged, the whole of it ended up inside the takeaway box.
        # y = 6.52 clears a gold rule the template master draws at y 6.08-6.39, x 1.50-5.76.
        # The first proof put this box at 5.72 and the goals ran straight through it. Template
        # furniture is part of the geometry whether or not this file drew it.
        add_text(
            slide,
            1.53,
            6.52,
            19.0,
            0.62,
            content["goals"],
            10.5,
            color=UW_GOLD,
            font=BODY_FONT,
            name="Header_Goals",
        )
    add_takeaway(slide, content, 1.50, 7.20, 21.0, 1.70)

    panels = content["panels"]
    for i, key in enumerate(str(n) for n in range(1, 10)):
        if key not in panels:
            continue
        set_cell(
            slide,
            figures,
            panels[key],
            COL_X[i % 3],
            ROW_Y[i // 3],
            COL_W,
            ROW_H,
            name=f"Cell{key}",
        )

    for i, key in enumerate(("strip1", "strip2", "strip3")):
        if key not in panels:
            continue
        set_cell(
            slide,
            figures,
            panels[key],
            COL_X[i],
            STRIP_Y,
            COL_W,
            STRIP_H,
            name=f"Strip{i + 1}",
        )

    add_footer(slide, content)
    return prs, figures


# ---------------------------------------------------------------------------------------------
# The factory board, v3. One builder, one geometry, no alternates.
# ---------------------------------------------------------------------------------------------
#
# Why this exists next to build_board rather than replacing it. build_board sets a uniform 3x3 of
# cells that each carry a heading, a big number, a figure and a paragraph. That grid is right for
# the top board, where nine findings really are nine of the same shape. It was wrong here, and
# the proof PNGs said so before anyone did: this board's charts were report figures drawn 8-16 in
# wide and scaled to about 0.3x inside a 2.12 in slot, so their axis type printed at 3-5 pt
# beside 18 pt prose. Roughly a seventh of the board was illegible chart.
#
# Jeffrey's rule, and it is the design now: a poster with three or four readable charts beats a
# poster with eight nobody can read, and 15-20 words under a chart beat 55.
#
# So the panels are no longer all the same shape. A chart panel is a heading, a chart drawn at
# the size it prints, and a caption of about 25 words. A statement card is a heading that IS the
# finding and about 40 words behind it, because its point was always a sentence rather than a
# picture. The rail is six numbers. Nothing was deleted to get here -- what came off the board
# moved into report 09, where it is read at reading distance.

# Every box on the board, in inches, on the 24 x 36 in sheet. Read down the page: the header band
# and its takeaway are add_header/add_takeaway's, unchanged; everything below is this.
#
# The columns are 0.90 and 23.10 at the edges with a 0.30 gutter, which is wider than v2's 1.50
# margin on purpose -- Jeffrey asked for less margin and more chart, and 1.20 in of returned
# width is most of a value label.
# The top board's column grid, not one of our own. Columns at 1.50 / 8.75 / 16.13, 6.35 in wide,
# 0.90 gutters, which is what the UW template lays out and what report 13's board uses. The two
# halves of the poster hang one above the other, so sharing a grid is worth more than any styling
# decision either board makes alone -- a reader's eye runs down a column across both.
COL_X = (1.50, 8.75, 16.13)
COL_W = 6.35
SPAN2 = 13.60          # two columns and the gutter between them
SPAN3 = 20.98          # the full measure

# Row tops and heights. The order is Jeffrey's: charts, the speedups table, charts, then the two
# text rows -- which follows the order the work happened in and puts a full-width band between the
# two chart rows so they do not read as one six-panel block.
# 8.40, not 8.05: the template master paints its purple band down to 8.24, and at 8.05 the first
# line of every heading in the top row printed behind it. Measured off the proof -- furniture this
# file did not draw is still geometry it has to live inside.
# Row A starts at 8.80, not 8.40. The template master paints its purple band to 8.24, and a
# heading landing 0.16 in under it reads as touching -- a band and a column of type want a margin
# between them, not a clearance.
_ROW = {'a': (8.70, 6.15), 'd': (15.00, 5.60), 'b': (20.75, 5.64),
        'e': (26.54, 3.82), 'c': (30.51, 4.30)}


def _across(row, w=COL_W, n=3):
    """The n column boxes of a row, or one box spanning the full measure."""
    y, h = _ROW[row]
    if n == 1:
        return ((COL_X[0], y, w, h),)
    return tuple((COL_X[i], y, w, h) for i in range(n))


def _split(row):
    """A row of two: two columns and a gutter, then the third column on its own."""
    y, h = _ROW[row]
    return ((COL_X[0], y, SPAN2, h), (COL_X[2], y, COL_W, h))


def _split_r(row):
    """The mirror: one column first, then the remaining two as a span."""
    y, h = _ROW[row]
    return ((COL_X[0], y, COL_W, h), (COL_X[1], y, SPAN2, h))


FACTORY_GEO = {
    "row_a": _across('a'),
    "row_d": _split_r('d'),
    "row_b": _across('b'),
    "row_e": _across('e'),
    "row_c": _split('c'),
}

# Which panel goes in which box, in the order the rows are drawn. Written down rather than derived
# from the numbering: the board file's order and the wall's order agreeing is a property worth
# being able to break, and this is the one place that changes when a panel moves.
FACTORY_PLAN = {
    "row_a": ("1", "8", "3"),                        # hardware · memory · the rail
    "row_d": ("2", "4"),                             # budgeting · what we made faster
    "row_b": ("5", "6", "7"),                        # early stopping · transfer · dashboard
    "row_e": ("9", "10", "strip1"),                  # identity · noise · limits
    "row_c": ("11", "strip2"),                        # the languages at two columns · sources
}

# What kind of block each panel is. A chart panel gets a figure; a rail is the purple stat block;
# a table is the speedups grid; everything else is a heading and prose. Only the two that cannot
# be told apart from the panel itself are listed -- a figure is a chart, and no figure is prose.
FACTORY_KIND = {"2": "list", "3": "rail", "4": "table", "8": "list", "11": "list2", "strip2": "refs"}

# Type. The top board sets section headers at 40 pt, but its headers are single words -- ABSTRACT,
# RESULTS -- and ours are whole statements. At 6.35 in a 40 pt sentence wraps to three lines and
# eats the chart, so panel headings run at 32 and the title takes the size back.
#
# The title also goes ALL CAPS in the theme's own major font, which is the Black weight and not
# the Normal weight we had been bolding. Headings stay sentence case: an all-caps SENTENCE at
# 32 pt is slower to read than the same words set in a heavier face, and these are sentences.
FACTORY_PT = {
    "title": 72.0,
    # One body size on the board: 17 pt, Jeffrey's number, everywhere prose or a list row
    # appears. Heads at 28 rather than 32 -- the 32/13 pairing put a canyon between a heading
    # and its own text, and the four points a heading gives up pay for the body's four.
    "panel_head": 28.0, "panel_cap": 17.0,
    "rail_head": 26.0, "rail_value": 30.0, "rail_label": 15.0,
    "table_head": 28.0, "table_row": 11.0, "table_lead": 17.0,
    "card_head": 28.0, "card_body": 17.0,
    "list_body": 17.0, "refs_body": 13.0, "footer": 12.0,
}


def add_rich_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float,
    *,
    color: str = INK,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float = 1.08,
    space_after: float = 0.0,
    margin: float = 0.04,
    name: str | None = None,
):
    """A textbox where `**...**` spans are bold, `*...*` italic, and newlines are paragraphs.

    add_text puts the whole string in one run, which is right for a heading and wrong for a
    caption. The board marks the phrase that carries the finding, and at two meters that mark is
    most of what a passer-by takes away -- rendering it as plain text throws away the only
    typographic decision the build sheet makes.

    Markdown emphasis is the only markup the board uses inside a panel body, so the parse is a
    split on the fence rather than anything general. Backticks come out entirely: a filename set
    in a proportional face beside prose reads as a typo, and the panels that name files name them
    in their provenance line, which does not print.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    set_cell_margins(tf, margin)
    # Paragraphs, because three panels are lists rather than prose -- five memory statements,
    # four language roles, six references. Newline in, paragraph out.
    for p, line in enumerate(text.split("\n")):
        paragraph = tf.paragraphs[0] if p == 0 else tf.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(space_after)
        paragraph.line_spacing = line_spacing
        # Both emphases, not just bold. The build sheet writes `**1.31x**` for a measurement and
        # `*slower*` for a distinction, and a parse that knew only the first printed the second
        # one's asterisks on the wall. Split keeps the delimiters so each run knows which it is.
        for chunk in re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", line.replace("`", "")):
            if not chunk:
                continue
            bold = chunk.startswith("**")
            italic = not bold and chunk.startswith("*")
            run = paragraph.add_run()
            run.text = chunk.strip("*")
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = rgb(color)
    return shape


def _panel_head(slide, x, y, w, text, pt, name, *, rule=True) -> float:
    """A heading and the hairline under it. Returns the y where the panel's content starts.

    This is what replaced the rounded rectangle. A box around every panel was the thing that read
    as belonging to a different poster than the one hanging above it -- the top board has no
    containers at all, just headings and whitespace over a shared grid. A 1 pt gold rule does the
    separating that a border was doing, and gives the chart back the 0.6 in the border and its
    padding were taking out of a 6.35 in column.
    """
    # Two lines at 28 pt is 68 pt; the box budgets for two because most headings take two.
    head_h = 0.98 if pt >= 26 else 0.62
    # Bottom-anchored. The box is sized for two lines because most headings take two, and a
    # one-line heading set from the top left half an inch of nothing between itself and its own
    # rule -- which reads as a panel whose heading has come adrift from its content.
    add_text(slide, x, y, w, head_h, text, pt, color=UW_PURPLE, font=HEADLINE_FONT,
             bold=True, valign=MSO_ANCHOR.BOTTOM, name=f"{name}_Heading")
    if rule:
        add_line(slide, x, y + head_h + 0.04, x + w, y + head_h + 0.04, UW_GOLD, 1.25)
    return y + head_h + 0.18


def add_chart_panel(slide, figures, panel, box, *, name: str, bleed: float = 0.0) -> None:
    """Heading, rule, chart, caption. The chart is drawn at exactly the size it lands.

    The caption is the panel's argument now. At 6.35 in a chart holds axis labels and value labels
    and nothing longer, so the sentences that used to sit inside the drawing -- figure 25's
    four-part key, figure 26's callout -- moved out here. That division is better than the one it
    replaced: the chart carries the picture and the caption carries the claim, and the caption is
    set in the same 18 pt as every other word on the board rather than in whatever a matplotlib
    text call happened to be scaled to.
    """
    x, y, w, h = box
    top = _panel_head(slide, x, y, w, panel.title, FACTORY_PT["panel_head"], name)
    cap_h = 1.60
    fig_h = (y + h - 0.14) - cap_h - 0.16 - top
    # Only the FIGURE bleeds -- into the page margin on the left and the gutter on the right --
    # so the panel still reads as a column while its chart gets the width Jeffrey asked for.
    # The heading and caption hold the column.
    figures.append(FigureSpec(_figure_source(panel.figure), x - bleed, top,
                              w + 2 * bleed, fig_h, name))
    add_rich_text(slide, x, top + fig_h + 0.16, w, cap_h, panel.body,
                  FACTORY_PT["panel_cap"], color=INK, space_after=6.0,
                  name=f"{name}_Caption")


def _figure_source(filename: str) -> Path:
    """The SVG the board names, or its PNG where only that has been rendered."""
    source = FIGURES / filename
    if not source.exists() and Path(filename).suffix.lower() == ".svg":
        source = FIGURES / f"{Path(filename).stem}.png"
    return source


def add_rail(slide, panel, box, *, name: str = "Rail") -> None:
    """The stat rail: the scale of the factory as gold numbers on purple.

    The one panel that keeps a fill, and it keeps it because the top board has the same device --
    its KEY DETAILS block is a purple panel with a gold heading and white text. Matching it is
    alignment; the rounded rectangles around everything else were not.

    Rows are value-label pairs split off the middots in the panel's own body, so the rail cannot
    say something the build sheet does not. The labels run two to three times longer than v1's --
    "pretraining runs" became a sentence that says what a pretraining run IS -- which is the whole
    of Jeffrey's note: a number nobody can read the units of is decoration.
    """
    x, y, w, h = box
    add_rect(slide, x, y, w, h, fill=UW_PURPLE, name=f"{name}_Box")
    add_text(slide, x + 0.30, y + 0.24, w - 0.60, 0.60, panel.title,
             FACTORY_PT["rail_head"], color=UW_GOLD, font=HEADLINE_FONT, bold=True,
             name=f"{name}_Heading")

    rows = []
    for part in panel.body.split(" · "):
        m = re.match(r"\s*\*\*(.+?)\*\*\s*(.*)", part.strip(), re.S)
        if m:
            rows.append((m.group(1).strip(), " ".join(m.group(2).split())))
    top = y + 0.98
    pitch = (h - 1.22) / max(len(rows), 1)
    for i, (value, label) in enumerate(rows):
        row_y = top + i * pitch
        add_text(slide, x + 0.28, row_y, 1.72, pitch, value, FACTORY_PT["rail_value"],
                 color=UW_GOLD, font=HEADLINE_FONT, bold=True, align=PP_ALIGN.RIGHT,
                 valign=MSO_ANCHOR.MIDDLE, name=f"{name}_Value{i + 1}")
        add_text(slide, x + 2.14, row_y, w - 2.44, pitch, label, FACTORY_PT["rail_label"],
                 color=PAPER, font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
                 name=f"{name}_Label{i + 1}")


def add_text_panel(slide, panel, box, *, name: str) -> None:
    """A heading that is the finding, and the sentences behind it. No figure, by choice.

    Statement cards and the strip blocks are the same thing now. They were different sizes when
    they lived in different-sized boxes; with the boxes gone they are both a heading and a
    paragraph in a 6.35 in column, and setting the strip two steps down was only ever compensating
    for a box that was too short.
    """
    x, y, w, h = box
    top = _panel_head(slide, x, y, w, panel.title, FACTORY_PT["card_head"], name)
    # 6 pt between paragraphs -- Jeffrey's hand edit on the printed board, alongside the
    # sentence-per-paragraph breaks the build sheet now records as bare '>' lines.
    add_rich_text(slide, x, top, w, (y + h - 0.06) - top, panel.body,
                  FACTORY_PT["card_body"], color=INK, space_after=6.0,
                  name=f"{name}_Body")


# The speedups table's three columns, as a fraction of one half-block's width. `change` takes most
# of it because that column is prose and the other two are a measurement and one word.
TABLE_COLS = (0.00, 0.440, 0.820)


def add_table_panel(slide, panel, box, *, name: str) -> None:
    """The speedups table: eight measured changes, in two blocks of four side by side.

    Two blocks rather than one column of eight. At 20.98 in a single eight-row table sets rows
    0.30 in apart with 15 in of white to the right of every entry, and the eye loses which
    measurement belongs to which change. Four rows in each half keeps the row band short enough to
    track and puts the type at 18 pt instead of the 13 it would need in one column.

    Rows come off the build sheet's own markdown table, so a corrected measurement is a one-line
    edit in report 12 rather than a value typed in here.
    """
    x, y, w, h = box
    top = _panel_head(slide, x, y, w, panel.title, FACTORY_PT["table_head"], name)
    add_rich_text(slide, x, top, w, 0.94, panel.body, FACTORY_PT["table_lead"],
                  color=INK, name=f"{name}_Lead")

    rows = list(panel.table)
    half = (len(rows) + 1) // 2
    block_w = (w - 0.60) / 2
    row_h = 0.235
    for b, chunk in enumerate((rows[:half], rows[half:])):
        bx = x + b * (block_w + 0.60)
        for c, head in enumerate(("what changed", "measured", "kept?")):
            add_text(slide, bx + TABLE_COLS[c] * block_w, top + 0.98, block_w * 0.30, 0.22,
                     head.upper(), 10.0, color=MUTED, font=SUBHEAD_FONT, bold=True,
                     name=f"{name}_H{b}{c}")
        for r, cells in enumerate(chunk):
            ry = top + 1.22 + r * row_h
            # A row whose other cells are empty is a SECTION, not data. The build sheet writes
            # `| **the data path** | | |` and gets the report's own grouping -- which is most of
            # why the twenty-change version reads better than a flat list of the same rows.
            if len(cells) > 1 and not any(c.strip() for c in cells[1:]):
                add_text(slide, bx, ry, block_w, row_h - 0.02,
                         cells[0].strip('*').upper(), FACTORY_PT["table_row"] - 1.0,
                         color=UW_PURPLE, font=SUBHEAD_FONT, bold=True,
                         valign=MSO_ANCHOR.MIDDLE, name=f"{name}_G{b}{r}")
                continue
            add_line(slide, bx, ry - 0.04, bx + block_w, ry - 0.04, LINE, 0.75)
            for c, text in enumerate(cells[:3]):
                cw = (TABLE_COLS[c + 1] if c + 1 < len(TABLE_COLS) else 1.0) - TABLE_COLS[c]
                # The verdict column is bold and colored: with twenty rows, the SHAPE of the
                # list -- mostly kept, a visible band of rejections -- is the panel's argument,
                # and a column of identical grey words hid it. Green kept, red rejected, blue
                # could-not; "no change" stays muted, being neither a win nor a rejection.
                if c == 2:
                    verdict = text.strip().lower()
                    colour = (KEPT_GREEN if verdict.startswith(("yes", "guard"))
                              else KEPT_RED if verdict == "no"
                              else KEPT_BLUE if verdict.startswith("can")
                              else MUTED)
                    add_text(slide, bx + TABLE_COLS[c] * block_w, ry, cw * block_w - 0.10,
                             row_h - 0.04, text, FACTORY_PT["table_row"], color=colour,
                             bold=True, name=f"{name}_R{b}{r}C{c}")
                    continue
                add_rich_text(slide, bx + TABLE_COLS[c] * block_w, ry, cw * block_w - 0.10,
                              row_h - 0.04, text, FACTORY_PT["table_row"],
                              color=INK, name=f"{name}_R{b}{r}C{c}")


def add_list_panel(slide, panel, box, *, name: str, row_break: bool = False) -> None:
    """A heading, a lead sentence, and the panel's table rows set as a list.

    Two panels want this shape and neither was working as prose. The memory panel is five
    statements each with its own evidence; the seventeen-languages panel is four roles, one per
    language group. Set as a paragraph, a reader at two meters has to find the boundaries
    themselves -- which is exactly what Jeffrey reported: the words were right and the panel did
    not land. A list puts the boundaries in the typography.

    The bold half of each row is the claim and the rest is what backs it, which is the build
    sheet's own two columns rendered on one line rather than in a grid: at 6.35 in a two-column
    table gives each cell three inches and every one of them wraps.
    """
    x, y, w, h = box
    top = _panel_head(slide, x, y, w, panel.title, FACTORY_PT["card_head"], name)
    lead_h = 0.68 if panel.body else 0.0
    if panel.body:
        add_rich_text(slide, x, top, w, lead_h, panel.body, FACTORY_PT["card_body"],
                      color=INK, name=f"{name}_Lead")
    # row_break once meant two lines per row; Jeffrey's hand edit on the printed board merged
    # them back to one paragraph per language with more air between rows, and he has now seen
    # both versions printed -- this records the later verdict.
    text = "\n".join(" ".join(c for c in cells if c.strip()) for cells in panel.table)
    # Jeffrey's hand-edit spacing, read off his printed-board file: 6 pt between list rows, and
    # the languages panel a step airier at 8 pt with 1.1 leading.
    add_rich_text(slide, x, top + lead_h + 0.08, w, (y + h - 0.06) - (top + lead_h + 0.08),
                  text, FACTORY_PT["list_body"], color=INK,
                  space_after=8.0 if row_break else 6.0,
                  line_spacing=1.1 if row_break else 1.02, name=f"{name}_List")


def add_refs_panel(slide, panel, box, *, name: str) -> None:
    """The reference list as ONE flowing block, hanging-indented, one citation per paragraph.

    It was a stack of per-reference boxes, each individually sized by a character-count
    estimator -- and the gate measures boxes, so the STACK ran off the panel twice without a
    single box overflowing. One box means one honest measurement, and the estimator retires.

    The hanging indent is raw pPr XML (marL positive, indent negative) because python-pptx does
    not surface paragraph indents; a wrapped citation returns under its own first character,
    which is how the top board's references read.
    """
    x, y, w, h = box
    top = _panel_head(slide, x, y, w, panel.title, 26.0, name)
    entries = "\n".join(e.strip() for e in panel.body.split(" · ") if e.strip())
    shape = add_rich_text(slide, x, top, w, (y + h - 0.06) - top, entries,
                          FACTORY_PT["refs_body"], color=INK, space_after=6.0,
                          name=f"{name}_Refs")
    for paragraph in shape.text_frame.paragraphs:
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set("marL", "182880")     # 0.20 in
        pPr.set("indent", "-182880")

def build_factory_poster(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    """The bottom board, v2. Five charts at column width, a rail, a table, six blocks of prose.

    The rows are drawn in the order FACTORY_PLAN lists them, which is the order the work happened
    in rather than the order the panels are numbered: the two chart rows are separated by the
    full-width speedups band so they do not read as one six-panel block.
    """
    prs, slide = open_template(PORTRAIT_TEMPLATE, 0)
    # The template LAYOUT ships three connectors for its sample blocks -- a full-width rule at
    # y 15.15 and column dividers at x 8.32 and 15.69 -- and they print through everything laid
    # on top. They were faint enough to survive four proofs before Jeffrey circled one. Deleted
    # from this presentation's in-memory layout only; the template file is untouched.
    for sh in list(slide.slide_layout.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.LINE:
            sh._element.getparent().remove(sh._element)
    figures: list[FigureSpec] = []
    # Names at 4.42 and the tagline at 5.05: on a research poster the team goes directly
    # under the title, which is where a reader looks for it. They were below the tagline and
    # above a gold rule the template draws, in the smallest type in the band.
    # The author line takes the subtitle's own face and size and centers in the space between
    # the template's gold bar and the bottom of the band -- Jeffrey's arrangement, again: it had
    # drifted to a smaller face hugging the bar.
    add_header(slide, content, title_pt=FACTORY_PT["title"], title_caps=True,
               title_font=TITLE_BLACK, subtitle_pt=19.0, author_pt=19.0)
    # Compressed from 1.70 in: the body starts at 8.30 to sit level with the top board's first
    # section, and this box is what was in the way.
    # The gold box carries GOALS now, not the takeaway. Jeffrey moved the 501/502/503/504 thesis
    # into the subtitle -- it is what the board is for, and it had been sitting in a box below the
    # band while a tagline held the slot above it. That empties the box, and Goals is the one
    # rubric item in the header with nowhere else to go.
    if content.get("goals"):
        add_takeaway(slide, {"takeaway": content["goals"]}, 1.50, 7.15, 21.0, 0.78,
                     label="GOALS")

    panels = content["panels"]
    for row, keys in FACTORY_PLAN.items():
        for box, key in zip(FACTORY_GEO[row], keys):
            panel = panels.get(key)
            if panel is None:
                continue
            name = f"Cell{key}"
            kindof = FACTORY_KIND.get(key, "chart" if panel.figure else "text")
            if kindof == "rail":
                add_rail(slide, panel, box, name=name)
            elif kindof == "table":
                add_table_panel(slide, panel, box, name=name)
            elif kindof == "list":
                add_list_panel(slide, panel, box, name=name)
            elif kindof == "list2":
                add_list_panel(slide, panel, box, name=name, row_break=True)
            elif kindof == "refs":
                add_refs_panel(slide, panel, box, name=name)
            elif kindof == "chart":
                add_chart_panel(slide, figures, panel, box, name=name,
                                bleed=0.55 if key == "1" else 0.0)
            else:
                add_text_panel(slide, panel, box, name=name)

    add_footer(slide, content)
    return prs, figures


def build_factory_poster_v3(content: dict, kind: str) -> tuple[Presentation, list[FigureSpec]]:
    """The bottom board on the template's SECOND page -- the top poster's own header.

    An A/B against build_factory_poster, at Jeffrey's request: the Second Page's purple band
    ends at 7.25 in against the First Page's 8.50, which frees 1.25 in. All of it goes to air --
    the same five rows at the same heights (every chart still lands 1:1 on its calibrated box),
    with the inter-row gaps grown from 0.18 in to ~0.43. Same panels, same build sheet, same
    figures; the only variables are the header and the breathing room.
    """
    prs, slide = open_template(PORTRAIT_TEMPLATE, 1)
    # Same stray connectors as the First Page layout, plus the gold-bar picture at 5.11-6.15,
    # which sits exactly where a three-line header needs its subtitle.
    for sh in list(slide.slide_layout.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.LINE or sh.name == "Picture 14":
            sh._element.getparent().remove(sh._element)
    figures: list[FigureSpec] = []

    # The header, laid straight onto the band: the layout's own wide wordmark sits at y 1.52,
    # then title, thesis, names -- the same three lines as v2, with the band's 7.25 in to
    # breathe in rather than 8.50 to fill.
    # 21.79 in wide -- Jeffrey's hand edit on the promoted board: at that width the title sets
    # on ONE line at 72 pt, which is what the shallower Second-Page band wants.
    add_text(slide, 1.50, 2.32, 21.79, 2.48,
             content["title"].upper(), FACTORY_PT["title"], color=PAPER, font=TITLE_BLACK,
             bold=True, valign=MSO_ANCHOR.MIDDLE, name="Header_Title")
    add_text(slide, 1.53, 4.92, 21.0, 0.95, content["subtitle"], 19.0, color=PAPER,
             font=SUBHEAD_FONT, name="Header_Subtitle")
    add_text(slide, 1.53, 6.10, 21.0, 0.60,
             content.get("author") or "CSED 504 · University of Washington · Summer 2026",
             19.0, color=UW_GOLD, font=SUBHEAD_FONT, valign=MSO_ANCHOR.MIDDLE,
             name="Header_Author")

    # The same rows at the same heights, re-based on the shallower band with uniform ~0.43 in
    # gaps -- computed from FACTORY_GEO so the two builders cannot drift apart in content.
    y3 = {'row_a': 7.55, 'row_d': 14.14, 'row_b': 20.17, 'row_e': 26.25, 'row_c': 30.50}
    panels = content["panels"]
    for row, keys in FACTORY_PLAN.items():
        for box, key in zip(FACTORY_GEO[row], keys):
            x, _, w, h = box
            box3 = (x, y3[row], w, h)
            panel = panels.get(key)
            if panel is None:
                continue
            name = f"Cell{key}"
            kindof = FACTORY_KIND.get(key, "chart" if panel.figure else "text")
            if kindof == "rail":
                add_rail(slide, panel, box3, name=name)
            elif kindof == "table":
                add_table_panel(slide, panel, box3, name=name)
            elif kindof == "list":
                add_list_panel(slide, panel, box3, name=name)
            elif kindof == "list2":
                add_list_panel(slide, panel, box3, name=name, row_break=True)
            elif kindof == "refs":
                add_refs_panel(slide, panel, box3, name=name)
            elif kindof == "chart":
                add_chart_panel(slide, figures, panel, box3, name=name,
                                bleed=0.55 if key == "1" else 0.0)
            else:
                add_text_panel(slide, panel, box3, name=name)

    add_footer(slide, content, footer_y=34.95)
    return prs, figures


# Per kind, because the two halves of the poster are no longer the same shape. The top board is
# nine findings of one shape and build_board's measured 3x3 is right for it; the bottom board is
# five charts, a rail, three statements and a strip, and only build_factory_poster knows that.
#
# The factory half keeps ONE builder on purpose. The three design-* alternates arrange semantic
# slots named before the grid was measured -- "question", "causal" -- so they can only ever show
# a subset of a board they predate, and a subset they choose themselves. Rendering three of them
# beside the real one produced four files a reader has to be told which of to print.
# A version suffix per half. The factory board is on its second redesign and Jeffrey keeps the
# previous one to compare against, so its files carry the version and do not collide.
VERSION = {"factory": "_v2", "yoruba": ""}

BUILDERS = {
    "factory": [("board", build_factory_poster),
                # The Second-Page A/B -- delete or promote on Jeffrey's verdict.
                ("board3", build_factory_poster_v3)],
    "yoruba": [
        # The board itself, set to report 12's measured grid. This is the poster.
        ("board", build_board),
        # Alternates inherited from PR #77, kept because they are cheap to render and useful to
        # look at side by side.
        ("design-01-classic", build_classic),
        ("design-02-narrative-spine", build_spine),
        ("design-04-editorial", build_editorial),
    ],
}


def save_base_presentation(prs: Presentation, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def source_aspect(path: Path) -> float:
    if path.suffix.lower() == ".svg":
        import xml.etree.ElementTree as ET

        root = ET.parse(path).getroot()
        view_box = root.attrib.get("viewBox")
        if view_box:
            _, _, width, height = [float(v) for v in view_box.replace(",", " ").split()]
            return width / height
        width = float(root.attrib["width"].replace("pt", "").replace("px", ""))
        height = float(root.attrib["height"].replace("pt", "").replace("px", ""))
        return width / height
    with Image.open(path) as image:
        return image.width / image.height


def contain_box(spec: FigureSpec) -> tuple[float, float, float, float]:
    aspect = source_aspect(spec.path)
    box_aspect = spec.w / spec.h
    if aspect >= box_aspect:
        width = spec.w
        height = width / aspect
        x = spec.x
        y = spec.y + (spec.h - height) / 2
    else:
        height = spec.h
        width = height * aspect
        x = spec.x + (spec.w - width) / 2
        y = spec.y
    return x, y, width, height


def _restore_reopened_picture_aspects(
    presentation, records, sx: float, sy: float
) -> None:
    """Correct the lazy PageSetup transform after reopening the file."""

    for record in records:
        if (
            record["container_kind"] != "slide"
            or not record["shape_name"].startswith("Embedded_")
        ):
            continue
        slide = presentation.Slides(record["container_index"])
        shape = slide.Shapes.Item(record["shape_name"])
        source = ROOT.parent / Path(shape.AlternativeText)
        aspect = source_aspect(source)
        target_left = record["left"] * sx
        target_top = record["top"] * sy
        target_width = record["width"] * sx
        target_height = record["height"] * sy
        if aspect >= target_width / target_height:
            new_width = target_width
            new_height = new_width / aspect
            new_left = target_left
            new_top = target_top + (target_height - new_height) / 2
        else:
            new_height = target_height
            new_width = new_height * aspect
            new_left = target_left + (target_width - new_width) / 2
            new_top = target_top
        try:
            shape.LockAspectRatio = 0
        except Exception:
            pass
        shape.Width = max(new_width, 0.1)
        shape.Height = max(new_height, 0.1)
        shape.Left = new_left
        shape.Top = new_top


def reflow_for_print(presentation, *, landscape: bool):
    """Convert 2:3 template geometry to the project's 3:4 print target.

    PowerPoint reflows shape boxes when PageSetup changes. We let it perform
    that axis-specific layout expansion, then restore picture aspect ratios
    and scale typography for the larger physical sheet.
    """

    if landscape:
        sx, sy = 4 / 3, 1.5
        font_scale = 1.5
        target_width, target_height = 48 * 72, 36 * 72
    else:
        sx, sy = 1.5, 4 / 3
        font_scale = 1.6
        target_width, target_height = 36 * 72, 48 * 72

    master = presentation.SlideMaster
    containers = [("master", 0, master)]
    for layout_index in range(1, master.CustomLayouts.Count + 1):
        containers.append(
            ("layout", layout_index, master.CustomLayouts(layout_index))
        )
    for slide_index in range(1, presentation.Slides.Count + 1):
        containers.append(("slide", slide_index, presentation.Slides(slide_index)))

    records = []
    for container_kind, container_index, container in containers:
        for shape_index in range(1, container.Shapes.Count + 1):
            shape = container.Shapes(shape_index)
            record = {
                "shape": shape,
                "container_kind": container_kind,
                "container_index": container_index,
                "shape_index": shape_index,
                "shape_name": shape.Name,
                "left": float(shape.Left),
                "top": float(shape.Top),
                "width": float(shape.Width),
                "height": float(shape.Height),
            }
            if int(shape.Type) == 13 and shape.Width > 0 and shape.Height > 0:
                record["picture_aspect"] = float(shape.Width) / float(shape.Height)
            try:
                if shape.HasTextFrame and shape.TextFrame2.HasText:
                    text_range = shape.TextFrame2.TextRange
                    if float(text_range.Font.Size) > 0:
                        record["font_size"] = float(text_range.Font.Size)
                    record["margins"] = (
                        float(shape.TextFrame2.MarginLeft),
                        float(shape.TextFrame2.MarginRight),
                        float(shape.TextFrame2.MarginTop),
                        float(shape.TextFrame2.MarginBottom),
                    )
            except Exception:
                pass
            try:
                if shape.Line.Visible:
                    record["line_weight"] = float(shape.Line.Weight)
            except Exception:
                pass
            records.append(record)

    presentation.PageSetup.SlideWidth = target_width
    presentation.PageSetup.SlideHeight = target_height

    for record in records:
        shape = record["shape"]
        target_left = record["left"] * sx
        target_top = record["top"] * sy
        target_box_width = record["width"] * sx
        target_box_height = record["height"] * sy
        if "picture_aspect" in record:
            continue
        else:
            new_left = target_left
            new_top = target_top
            new_width = target_box_width
            new_height = target_box_height
        try:
            shape.LockAspectRatio = 0
        except Exception:
            pass
        shape.Left = new_left
        shape.Top = new_top
        shape.Width = max(new_width, 0.1)
        shape.Height = max(new_height, 0.1)
        if "font_size" in record:
            try:
                shape.TextFrame2.TextRange.Font.Size = (
                    record["font_size"] * font_scale
                )
            except Exception:
                pass
        if "margins" in record:
            try:
                margins = record["margins"]
                shape.TextFrame2.MarginLeft = margins[0] * 4 / 3
                shape.TextFrame2.MarginRight = margins[1] * 4 / 3
                shape.TextFrame2.MarginTop = margins[2] * 4 / 3
                shape.TextFrame2.MarginBottom = margins[3] * 4 / 3
            except Exception:
                pass
        if "line_weight" in record:
            try:
                shape.Line.Weight = record["line_weight"] * 4 / 3
            except Exception:
                pass
    return records, sx, sy


def embed_figures_and_export(
    app,
    pptx_path: Path,
    figures: Iterable[FigureSpec],
    *,
    landscape: bool,
    reflow: bool = False,
) -> dict:
    """Place the figures, optionally rescale the sheet, export PDF + PNG, and report overflows.

    `reflow=False` is the default as of 17 August, and it is a size decision rather than a
    tidy-up. `reflow_for_print` blows the 24 x 36 template up to 36 x 48 and multiplies every
    font by 1.6 -- so every poster this module produced was 36 x 48, which is precisely the size
    #78 rejected, deleted eight files over, and added `assert_template_size` to prevent. That
    assertion passes anyway: it measures the *template* when it is opened, and the rescale
    happens afterwards on the way out. A guard on the input to a transform says nothing about
    its output.

    Report 12 settled the size by measuring the real template -- 24 x 36, body 18 pt, ~55 words
    a cell -- and the bottom board was rewritten around it. Building to that and then printing
    at 36 x 48 is how the two halves of this repository came to disagree about the board.
    """
    presentation = app.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
    slide = presentation.Slides(1)
    for index, spec in enumerate(figures, start=1):
        x, y, width, height = contain_box(spec)
        shape = slide.Shapes.AddPicture(
            str(spec.path.resolve()),
            False,
            True,
            x * 72,
            y * 72,
            width * 72,
            height * 72,
        )
        shape.Name = f"Embedded_{index:02d}_{spec.name}"
        shape.AlternativeText = str(spec.path.relative_to(ROOT.parent))
    if reflow:
        reflow_records, sx, sy = reflow_for_print(
            presentation, landscape=landscape
        )
        # PowerPoint applies its page-size transform lazily on the first save.
        # Reopen before reasserting picture aspect ratios.
        presentation.Save()
        presentation.Close()
        presentation = app.Presentations.Open(
            str(pptx_path.resolve()), WithWindow=False
        )
        _restore_reopened_picture_aspects(presentation, reflow_records, sx, sy)
    presentation.Save()
    slide = presentation.Slides(1)

    pdf_path = pptx_path.with_suffix(".pdf")
    presentation.SaveAs(str(pdf_path.resolve()), 32)
    preview_path = pptx_path.with_name(f"{pptx_path.stem}-preview.png")
    # Export at the SLIDE'S aspect, not a hardcoded pair. 1800 x 2400 is 0.750 wide-to-tall and a
    # 24 x 36 board is 0.667, so every preview was stretched ~12% horizontally -- the PDFs were
    # always right, but the PNG is what gets pasted into an email and reviewed, and it was quietly
    # showing a poster nobody was going to print. Found by cropping a preview against shape
    # coordinates read out of the .pptx and having the two disagree.
    page_w = float(presentation.PageSetup.SlideWidth)
    page_h = float(presentation.PageSetup.SlideHeight)
    long_edge = 2400
    if page_w >= page_h:
        px_w, px_h = long_edge, max(1, round(long_edge * page_h / page_w))
    else:
        px_h, px_w = long_edge, max(1, round(long_edge * page_w / page_h))
    slide.Export(str(preview_path.resolve()), "PNG", px_w, px_h)

    overflows = []
    for shape in slide.Shapes:
        try:
            if shape.HasTextFrame and shape.TextFrame2.HasText:
                bound_h = float(shape.TextFrame2.TextRange.BoundHeight)
                available_h = float(shape.Height)
                if bound_h > available_h + 2:
                    overflows.append(
                        {
                            "shape": shape.Name,
                            "bound_height_pt": round(bound_h, 1),
                            "box_height_pt": round(available_h, 1),
                        }
                    )
        except Exception:
            continue
    # Measured on the way OUT, which is the half `assert_template_size` cannot see.
    page = (
        round(float(presentation.PageSetup.SlideWidth) / 72, 2),
        round(float(presentation.PageSetup.SlideHeight) / 72, 2),
    )
    presentation.Close()
    return {
        "pptx": pptx_path.name,
        "pdf": pdf_path.name,
        "preview": preview_path.name,
        "page_inches": list(page),
        "text_overflows": overflows,
    }


def render_pdf_first_page(pdf_path: Path, png_path: Path) -> None:
    document = fitz.open(pdf_path)
    page = document[0]
    pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
    pix.save(png_path)
    document.close()


def build_contact_sheet(previews: list[Path], out_path: Path) -> None:
    thumb_w, thumb_h = 560, 840
    margin, label_h = 24, 42
    canvas = Image.new(
        "RGB",
        (2 * thumb_w + 3 * margin, 4 * (thumb_h + label_h) + 5 * margin),
        "white",
    )
    draw = ImageDraw.Draw(canvas)
    try:
        font = ImageFont.truetype("arial.ttf", 22)
    except OSError:
        font = ImageFont.load_default()
    for index, path in enumerate(previews):
        row, col = divmod(index, 2)
        with Image.open(path) as image:
            image = image.convert("RGB")
            image.thumbnail((thumb_w, thumb_h))
            x = margin + col * (thumb_w + margin)
            y = margin + row * (thumb_h + label_h + margin)
            frame_x = x + (thumb_w - image.width) // 2
            frame_y = y + (thumb_h - image.height) // 2
            canvas.paste(image, (frame_x, frame_y))
            draw.rectangle(
                (frame_x, frame_y, frame_x + image.width, frame_y + image.height),
                outline=(210, 205, 216),
                width=2,
            )
            draw.text((x, y + thumb_h + 8), path.stem.replace("-preview", ""), fill=(50, 0, 110), font=font)
    canvas.save(out_path, quality=95)


def write_readme() -> None:
    text = """# A2-NLP poster concepts

Four visual design systems are provided. Each design has a separate Yoruba findings poster and
a separate Model Factory poster.

| Design | Format | Intent |
|---|---|---|
| 01 Classic | 36 × 48 in portrait | Familiar research-poster hierarchy and strong figure grid |
| 02 Narrative spine | 36 × 48 in portrait | Central evidence chain / workflow with supporting proof |
| 03 Evidence grid | 48 × 36 in landscape | Wide presentation format; fast left-to-right scan |
| 04 Editorial | 36 × 48 in portrait | Fewer, larger statements for distance readability |

Each poster is supplied as an editable `.pptx`, print-oriented `.pdf`, and `.png` preview. Existing
repository figures are embedded as SVG where available so charts remain vector in PowerPoint/PDF.
The supplied 2:3/3:2 templates are reflowed—without stretching figures or logos—to the project's
3:4 physical-board specification.

Typography note: the supplied Encode Sans, Open Sans and Uni Sans files do not collectively cover
all precomposed and combining Yoruba characters. The current English-language poster copy does not
set Yoruba-script examples. If examples are added, use a Yoruba-complete font such as Noto Sans and
preflight `Ẹ/ẹ`, `Ọ/ọ`, `Ṣ/ṣ` and combining tone marks before printing.

The copy follows reports 11–13 and the current claims audit. In particular:

- SIB-200 says **ahead**, not **beats**, because bootstrap intervals overlap by 0.004.
- MasakhaNER rates were not selected under the same held-out protocol.
- The unsupported matched-compute pretraining mean penalty is not presented as a finding.
- The variability result and the matched-compute downstream vocabulary swap are kept distinct.

`build_posters.py` regenerates all files. It requires `python-pptx`, `pywin32`, Pillow and PyMuPDF,
and uses installed Microsoft PowerPoint to embed SVGs and export PDFs/previews.
"""
    (OUTPUTS / "README.md").write_text(text, encoding="utf-8")


def main() -> None:
    # The board check runs first and blocks. Rendering a cell whose prose has outgrown its column
    # produces a poster that looks finished and is not, and the overflow report at the end of this
    # function only measures boxes that were already drawn -- so a word budget has to be checked
    # before anything is placed, not after.
    problems = check_board()
    if problems:
        raise SystemExit(
            "board is not settable; fix these before rendering:\n  "
            + "\n  ".join(problems)
        )
    for note in long_panels() + stale_counts():
        print(f"note: {note}")

    OUTPUTS.mkdir(parents=True, exist_ok=True)

    jobs: list[tuple[Path, list[FigureSpec], bool]] = []
    for kind in ("yoruba-findings", "model-factory"):
        which = "yoruba" if kind == "yoruba-findings" else "factory"
        for design_name, builder in BUILDERS[which]:
            content = content_from_board(which)
            prs, figures = builder(content, which)
            # The v3 experiment keeps the family name rather than inheriting the design key.
            if design_name == "board3":
                path = OUTPUTS / "board-model-factory_v3.pptx"
            else:
                path = OUTPUTS / f"{design_name}-{kind}{VERSION[which]}.pptx"
            # Only what this run writes. The glob that used to sit above this loop deleted
            # every pptx, pdf and preview in the directory, including a saved earlier version
            # somebody had put there on purpose.
            # Windows holds an exclusive lock on an open .pptx, and this file is the thing
            # somebody is looking at while the build runs -- so the failure is normal, frequent,
            # and worth naming rather than raising a PermissionError from pathlib.
            for stale in (path, path.with_suffix(".pdf"),
                          path.with_name(path.stem + "-preview.png")):
                if not stale.exists():
                    continue
                try:
                    stale.unlink()
                except PermissionError:
                    raise SystemExit(
                        f"{stale.name} is open in another program -- close it and run again. "
                        f"Windows will not let the build replace a file PowerPoint has open, "
                        f"and half-writing it would be worse than stopping."
                    ) from None
            save_base_presentation(prs, path)
            landscape = prs.slide_width > prs.slide_height
            jobs.append((path, figures, landscape))

    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True
    reports = []
    try:
        for path, figures, landscape in jobs:
            reports.append(
                embed_figures_and_export(
                    app, path, figures, landscape=landscape
                )
            )
    finally:
        app.Quit()

    # The guard #78 meant to have. `assert_template_size` checks the template as it is opened;
    # this checks the sheet that actually came out, which is where the 36 x 48 was hiding.
    wrong_size = [r for r in reports if tuple(r["page_inches"]) != (24.0, 36.0)]
    if wrong_size:
        raise SystemExit(
            "exported at a size nobody can print:\n  "
            + "\n  ".join(f"{r['pptx']}: {r['page_inches'][0]} x {r['page_inches'][1]} in"
                          for r in wrong_size)
        )

    # And the overflow report BLOCKS. It spent its whole life as a field in a JSON file nobody
    # read, which is how it went unnoticed that autofit had made it incapable of firing. A gate
    # that only reports is a gate you find out about afterwards. Everything is written by this
    # point, so the outputs are on disk to look at -- the exit code is what says not to print them.
    over = [(r, o) for r in reports for o in r["text_overflows"]]
    if over:
        raise SystemExit(
            f"{len(over)} text box(es) overflow -- do not print these:\n  "
            + "\n  ".join(
                f"{r['pptx']}: {o['shape']} needs {o['bound_height_pt']} pt "
                f"in {o['box_height_pt']} pt"
                for r, o in over
            )
        )

    preview_paths = [OUTPUTS / report["preview"] for report in reports]
    build_contact_sheet(preview_paths, OUTPUTS / "poster-designs-contact-sheet.png")
    (OUTPUTS / "layout-validation.json").write_text(
        json.dumps(reports, indent=2), encoding="utf-8"
    )
    write_readme()

    stray = OUTPUTS / "_template_previews"
    if stray.exists():
        shutil.rmtree(stray)

    print(json.dumps(reports, indent=2))


if __name__ == "__main__":
    main()
